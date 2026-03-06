"""
timeline_slide_generator.py

Generate a single PowerPoint timeline slide from:
- close date
- process type (standard | accelerated)
- theme (light | dark)

Install:
  pip install python-pptx python-dateutil Pillow

Run:
  python timeline_slide_generator.py --close 2026-11-07 --process standard --theme dark --out timeline_slide.pptx

  # Custom subtitle:
  python timeline_slide_generator.py --close 2026-11-07 --process standard --theme dark \\
      --subtitle "Begin process on April 1st, launch to market early June, with expected close in November" \\
      --out timeline_slide.pptx

Logo files (optional — place in same directory as script):
  logo_dark.png  -> white/light logo for dark theme
  logo_light.png -> coloured logo for light theme

=== TEMPLATE CALIBRATION NOTES ===
All offsets derived from 4 real Carlsquare slides (Eventbase, Eyrus, Ayla, Guru).

STANDARD anchors (Eventbase 11/07 and Eyrus 09/10 consensus):
  launch_to_market: -22w  (both examples: -21.7w and -22.0w)
  ioi_due:          -14w  (both: -14.4w)
  loi_due:          -11w  (avg -10.9w)
  process_start:    -31w  (both: ~-31.5w)

ACCELERATED anchors (Guru 08/28):
  launch_to_market: -27w  (-26.6w)
  ioi_due:          -19w  (-19.3w)
  loi_due:          -12w  (-12.3w)
  process_start:    -34w  (-34.1w)

Standard task bars (confirmed from Eventbase visual):
  1a: -31w → -27w   (4w, ~1 month)
  1c: -31w → -23w   (8w, 2 months, overlaps 1a)
  1d: -31w → -22w   (9w, ends at launch, overlaps 1a/1c)
  1e: -29w → -21w   (8w, MP deck overlaps with 1a-1d)
  1f: -25w → -22w   (3w, fireside chats end at launch)
  2a: -22w → -14w   (8w, launch→IOI, "outreach and discussions" merged)
  2b: -20w → -14w   (6w, Distribute CIM)
  2d: -14w → -11w   (3w, mgmt presentations, starts at IOI)
  3a: -11w → -3w    (8w, due diligence, starts at LOI)
  3b: -8w  → -2w    (6w, negotiate PA, overlaps 3a)

Accelerated task bars (from Guru visual):
  1a: -34w → -30w   (4w)
  1c: -34w → -28w   (6w)
  1d: -34w → -27w   (7w, ends at launch)
  1e: -30w → -24w   (6w, MP deck)
  2a: -29w → -23w   (6w, discussions start 2w pre-launch)
  2b: -27w → -24w   (3w, outreach starts at launch)
  2c: -27w → -20w   (7w, Distribute CIM)
  2e: -19w → -12w   (7w, mgmt presentations, IOI→LOI)
  3a: -12w → -3w    (9w, due diligence starts at LOI)
  3b: -8w  → -1w    (7w, negotiate PA, overlaps 3a)
"""

from __future__ import annotations

import argparse
import os
from dataclasses import dataclass
from datetime import date, timedelta
from typing import Dict, List, Optional, Tuple

from dateutil.parser import isoparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ─────────────────────────────
# Templates
# Offsets calibrated from 4 real Carlsquare slides (Eventbase, Eyrus, Ayla, Guru)
# All offsets are in weeks relative to the close date (negative = before close)
# ─────────────────────────────
TEMPLATES = {
    "standard": {
        # ── Milestone anchors ──────────────────────────────────
        # Consensus from Eventbase (close 11/07) and Eyrus (close 09/10):
        #   launch_to_market: avg -21.9w → -22w
        #   ioi_due:          avg -14.4w → -14w
        #   loi_due:          avg -10.9w → -11w
        "anchors": {
            "close":            0,
            "loi_due":        -11,
            "ioi_due":        -14,
            "launch_to_market": -22,
        },
        "phases": [
            {
                "phase_id": "prep",
                "phase_label": "1. Preparation of documents and data room",
                "rows": [
                    # 1a: ~1 month, starts at process start (-31w), ends ~4w later
                    {"row_id": "1a", "label": "a. Preparation of investor / buyer list",
                     "include": True,
                     "start_offset_weeks": -31, "end_offset_weeks": -27},

                    # 1b: Quality of Earnings — optional/off by default in standard
                    {"row_id": "1b", "label": "b. Quality of Earnings",
                     "include": False,
                     "start_offset_weeks": -31, "end_offset_weeks": -25},

                    # 1c: CIM/Teaser ~2 months, overlaps 1a
                    {"row_id": "1c", "label": "c. Preparation of CIM and Teaser",
                     "include": True,
                     "start_offset_weeks": -31, "end_offset_weeks": -23},

                    # 1d: Data room ends at launch (-22w)
                    {"row_id": "1d", "label": "d. Data room (data collection & preparation)",
                     "include": True,
                     "start_offset_weeks": -31, "end_offset_weeks": -22},

                    # 1e: MP/fireside deck, starts ~2w after 1a, overlaps 1c/1d
                    {"row_id": "1e", "label": "e. Preparation of MP and fireside chat deck",
                     "include": True,
                     "start_offset_weeks": -29, "end_offset_weeks": -21},

                    # 1f: Fireside chats, ends at launch
                    {"row_id": "1f", "label": "f. Fireside chats",
                     "include": True,
                     "start_offset_weeks": -25, "end_offset_weeks": -22},
                ],
            },
            {
                "phase_id": "marketing",
                "phase_label": "2. Marketing / investor approach",
                "rows": [
                    # 2a: Outreach + discussions merged — spans launch → IOI (8w)
                    {"row_id": "2a", "label": "a. Investor / buyer outreach and discussions",
                     "include": True,
                     "start_offset_weeks": -22, "end_offset_weeks": -14},

                    # 2b: Distribute CIM, starts ~2w after launch, ends at IOI
                    {"row_id": "2b", "label": "b. Distribute CIM",
                     "include": True,
                     "start_offset_weeks": -20, "end_offset_weeks": -14},

                    # 2c: IOI milestone
                    {"row_id": "2c", "label": "c. Indications of Interest (IOI) due",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "ioi_due"},

                    # 2d: In-person mgmt presentations, starts at IOI, ends at LOI
                    {"row_id": "2d", "label": "d. In-person management presentations",
                     "include": True,
                     "start_offset_weeks": -14, "end_offset_weeks": -11},

                    # 2e: LOI milestone
                    {"row_id": "2e", "label": "e. Letters of Intent (LOI) due",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "loi_due"},
                ],
            },
            {
                "phase_id": "execution",
                "phase_label": "3. Execution",
                "rows": [
                    # 3a: Due diligence, starts at LOI, ~8 weeks
                    {"row_id": "3a", "label": "a. Due diligence",
                     "include": True,
                     "start_offset_weeks": -11, "end_offset_weeks": -3},

                    # 3b: Negotiate PA, starts ~3w after LOI, overlaps 3a, ends ~1w before close
                    {"row_id": "3b", "label": "b. Negotiate Purchase Agreement",
                     "include": True,
                     "start_offset_weeks": -8, "end_offset_weeks": -2},

                    # 3c: Close milestone
                    {"row_id": "3c", "label": "c. Sign and close",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "close"},
                ],
            },
        ],
    },

    "accelerated": {
        # ── Milestone anchors ──────────────────────────────────
        # From Guru (close 08/28/2026):
        #   launch_to_market: -26.6w → -27w
        #   ioi_due:          -19.3w → -19w
        #   loi_due:          -12.3w → -12w
        "anchors": {
            "close":            0,
            "loi_due":        -12,
            "ioi_due":        -19,
            "launch_to_market": -27,
        },
        "phases": [
            {
                "phase_id": "prep",
                "phase_label": "1. Preparation of documents and data room",
                "rows": [
                    # 1a: ~4 weeks, starts at process start (-34w)
                    {"row_id": "1a", "label": "a. Preparation of buyer list",
                     "include": True,
                     "start_offset_weeks": -34, "end_offset_weeks": -30},

                    # 1b: Quality of Earnings — shown in Guru (include True here)
                    {"row_id": "1b", "label": "b. Quality of Earnings",
                     "include": True,
                     "start_offset_weeks": -34, "end_offset_weeks": -27},

                    # 1c: CIM/Teaser ~6 weeks
                    {"row_id": "1c", "label": "c. Preparation of CIM and Teaser",
                     "include": True,
                     "start_offset_weeks": -34, "end_offset_weeks": -28},

                    # 1d: Data room ends at launch (-27w)
                    {"row_id": "1d", "label": "d. Data room (data collection & preparation)",
                     "include": True,
                     "start_offset_weeks": -34, "end_offset_weeks": -27},

                    # 1e: MP/fireside deck, ~6 weeks
                    {"row_id": "1e", "label": "e. Preparation of MP and fireside chat deck",
                     "include": True,
                     "start_offset_weeks": -30, "end_offset_weeks": -24},
                ],
            },
            {
                "phase_id": "marketing",
                "phase_label": "2. Marketing / investor approach",
                "rows": [
                    # 2a: Investor discussions — starts 2w BEFORE launch (Guru shows this)
                    {"row_id": "2a", "label": "a. Investor / buyer discussions",
                     "include": True,
                     "start_offset_weeks": -29, "end_offset_weeks": -23},

                    # 2b: Buyer outreach — starts at launch, ~3 weeks
                    {"row_id": "2b", "label": "b. Buyer outreach with teaser and NDAs",
                     "include": True,
                     "start_offset_weeks": -27, "end_offset_weeks": -24},

                    # 2c: Distribute CIM — starts at launch, ends ~7w later
                    {"row_id": "2c", "label": "c. Distribute CIM",
                     "include": True,
                     "start_offset_weeks": -27, "end_offset_weeks": -20},

                    # 2d: IOI milestone
                    {"row_id": "2d", "label": "d. Indications of Interest (IOI) due",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "ioi_due"},

                    # 2e: Mgmt presentations — IOI to LOI (~7w)
                    {"row_id": "2e", "label": "e. Management presentations",
                     "include": True,
                     "start_offset_weeks": -19, "end_offset_weeks": -12},

                    # 2f: LOI milestone
                    {"row_id": "2f", "label": "f. Letter of Intent (LOI) due",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "loi_due"},
                ],
            },
            {
                "phase_id": "execution",
                "phase_label": "3. Execution",
                "rows": [
                    # 3a: Due diligence, starts at LOI, ~9 weeks
                    {"row_id": "3a", "label": "a. Due diligence",
                     "include": True,
                     "start_offset_weeks": -12, "end_offset_weeks": -3},

                    # 3b: Negotiate PA, starts ~4w after LOI, overlaps 3a, ends ~1w before close
                    {"row_id": "3b", "label": "b. Negotiate Purchase Agreement",
                     "include": True,
                     "start_offset_weeks": -8, "end_offset_weeks": -1},

                    # 3c: Close milestone
                    {"row_id": "3c", "label": "c. Sign and close",
                     "include": True,
                     "type": "milestone_row", "milestone_key": "close"},
                ],
            },
        ],
    },
}


# ─────────────────────────────
# Themes
# ─────────────────────────────
THEMES = {
    "light": {
        "bg":           RGBColor(255, 255, 255),
        "text":         RGBColor(20,  20,  20),
        "muted_text":   RGBColor(80,  80,  80),
        "grid":         RGBColor(200, 200, 200),
        "phase_fill":   RGBColor(242, 242, 242),
        "month_strip":  RGBColor(18,  33,  61),
        "month_text":   RGBColor(255, 255, 255),
        "bar":          RGBColor(0,   154, 143),
        "launch_line":  RGBColor(80,  80,  200),
        "milestone":    RGBColor(92,  45,  145),
        "footer_text":  RGBColor(100, 100, 100),
        "top_bar":      RGBColor(18,  33,  61),
        "logo_file":    "logo_light.png",
    },
    "dark": {
        "bg":           RGBColor(10,  24,  48),
        "text":         RGBColor(255, 255, 255),
        "muted_text":   RGBColor(200, 210, 225),
        "grid":         RGBColor(60,  85,  120),
        "phase_fill":   RGBColor(20,  45,  80),
        "month_strip":  RGBColor(20,  45,  80),
        "month_text":   RGBColor(255, 255, 255),
        "bar":          RGBColor(114, 66,  244),
        "launch_line":  RGBColor(210, 210, 210),
        "milestone":    RGBColor(245, 90,  160),
        "footer_text":  RGBColor(180, 195, 215),
        "top_bar":      RGBColor(200, 210, 225),
        "logo_file":    "logo_dark.png",
    },
}


# ─────────────────────────────
# Date helpers
# ─────────────────────────────
def first_day_of_month(d: date) -> date:
    return date(d.year, d.month, 1)


def add_months(d: date, months: int) -> date:
    y = d.year + (d.month - 1 + months) // 12
    m = (d.month - 1 + months) % 12 + 1
    return date(y, m, 1)


def last_day_of_month(d: date) -> date:
    return add_months(first_day_of_month(d), 1) - timedelta(days=1)


def month_range(axis_start: date, axis_end: date) -> List[Tuple[date, date, str]]:
    months, cur = [], first_day_of_month(axis_start)
    end_month   = first_day_of_month(axis_end)
    while cur <= end_month:
        months.append((cur, last_day_of_month(cur), cur.strftime("%b")))
        cur = add_months(cur, 1)
    return months


# ─────────────────────────────
# Data model
# ─────────────────────────────
@dataclass
class DisplayRow:
    kind:          str
    phase_label:   Optional[str]
    row_id:        Optional[str]
    label:         str
    start:         Optional[date] = None
    end:           Optional[date] = None
    milestone_key: Optional[str]  = None


def compute_schedule(close_date: date,
                     tmpl: Dict) -> Tuple[List[DisplayRow], Dict[str, date]]:
    milestone_dates = {k: close_date + timedelta(days=wk * 7)
                       for k, wk in tmpl["anchors"].items()}
    rows: List[DisplayRow] = []
    for phase in tmpl["phases"]:
        rows.append(DisplayRow(kind="phase", phase_label=phase["phase_label"],
                               row_id=None, label=phase["phase_label"]))
        for r in phase["rows"]:
            if not r.get("include", True):
                continue
            if r.get("type") == "milestone_row":
                rows.append(DisplayRow(
                    kind="milestone", phase_label=phase["phase_label"],
                    row_id=r["row_id"], label=r["label"],
                    milestone_key=r["milestone_key"],
                ))
            else:
                rows.append(DisplayRow(
                    kind="task", phase_label=phase["phase_label"],
                    row_id=r["row_id"], label=r["label"],
                    start=close_date + timedelta(days=int(r["start_offset_weeks"]) * 7),
                    end  =close_date + timedelta(days=int(r["end_offset_weeks"])   * 7),
                ))
    return rows, milestone_dates


# ─────────────────────────────
# Shape / text helpers
# ─────────────────────────────
def solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def set_shape_line(shape, rgb: RGBColor, width_pt: float = 0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_textbox(slide, x, y, w, h, text,
                font_size=11, bold=False,
                color=RGBColor(0, 0, 0),
                align=PP_ALIGN.LEFT,
                italic=False,
                font_name="Roboto",
                vcenter=False):
    from pptx.enum.text import MSO_ANCHOR
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.clear()
    if vcenter:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    p.alignment        = align
    tf.word_wrap       = True
    return box


def dashed_vline(slide, x, y0, y1, color: RGBColor, dash=0.10, gap=0.07):
    """Simulate a dashed vertical line with stacked thin rectangles."""
    y = y0
    while y < y1:
        seg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.012), Inches(min(y + dash, y1) - y)
        )
        solid_fill(seg, color)
        y += dash + gap


# ─────────────────────────────
# Main render
# ─────────────────────────────
def render_timeline_slide(
    close_date:  date,
    process:     str,
    theme_name:  str,
    out_path:    str,
    subtitle:    str = "",
    top_label:   str = "[Insert section label here]",
):
    if process    not in TEMPLATES: raise ValueError(f"Unknown process: {process}")
    if theme_name not in THEMES:    raise ValueError(f"Unknown theme: {theme_name}")

    tmpl  = TEMPLATES[process]
    theme = THEMES[theme_name]

    rows, milestone_dates = compute_schedule(close_date, tmpl)

    # Milestone row_id lookup (which row each milestone sits on)
    milestone_row_map = {}
    for r in rows:
        if r.kind == "milestone" and r.milestone_key:
            milestone_row_map[r.milestone_key] = r.row_id

    # ── Axis bounds ──────────────────────────────────────────
    all_dates = (
        [r.start for r in rows if r.start] +
        [r.end   for r in rows if r.end]   +
        list(milestone_dates.values())
    )
    axis_start = first_day_of_month(min(all_dates))
    axis_end   = last_day_of_month(max(all_dates))
    months     = month_range(axis_start, axis_end)
    total_days = (axis_end - axis_start).days + 1

    # ── Slide: 4:3  (10" × 7.5") ─────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(10.0)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    SW, SH = 10.0, 7.5

    # Background
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 Inches(0), Inches(0), Inches(SW), Inches(SH))
    solid_fill(bg, theme["bg"])

    # ── Layout constants ──────────────────────────────────────
    M_left   = 0.28
    M_right  = 0.18
    Left_w   = 3.10
    Footer_y = 7.12
    Footer_h = 0.28

    TopLabel_y   = 0.15
    TopLabel_h   = 0.22
    Title_y      = TopLabel_y + TopLabel_h + 0.04
    Subtitle_y   = Title_y + 0.46
    MonthRow_y   = Subtitle_y + 0.36
    MonthRow_h   = 0.36
    Chart_bottom = Footer_y - 0.06

    Grid_x0 = M_left + Left_w
    Grid_x1 = SW - M_right
    Grid_w  = Grid_x1 - Grid_x0
    Grid_y0 = MonthRow_y + MonthRow_h
    Grid_y1 = Chart_bottom

    def x_of(d: date) -> float:
        return Grid_x0 + (d - axis_start).days / total_days * Grid_w

    # ── Top label ─────────────────────────────────────────────
    add_textbox(slide, M_left, TopLabel_y,
                SW - M_left - M_right, TopLabel_h,
                top_label,
                font_size=8, bold=False, color=theme["top_bar"])

    # ── Title ─────────────────────────────────────────────────
    add_textbox(slide, M_left, Title_y,
                SW - M_left - M_right, 0.44,
                "Timeline optimized for successful outcome",
                font_size=20, bold=True, color=theme["text"])

    # ── Subtitle ─────────────────────────────────────────────
    launch = milestone_dates["launch_to_market"]
    if not subtitle:
        subtitle = (
            f"Launch to market {launch.strftime('%b %d, %Y')}  •  "
            f"Expected close {close_date.strftime('%b %d, %Y')}"
        )
    add_textbox(slide, M_left, Subtitle_y,
                SW - M_left - M_right, 0.30,
                subtitle,
                font_size=10, bold=True, color=theme["text"])

    # ── Month header strip ────────────────────────────────────
    for i, (m_start, m_end, m_label) in enumerate(months):
        x0 = x_of(m_start)
        x1 = x_of(m_end + timedelta(days=1))
        cell_w = x1 - x0

        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                       Inches(x0), Inches(MonthRow_y),
                                       Inches(cell_w), Inches(MonthRow_h))
        solid_fill(cell, theme["month_strip"])

        if i < len(months) - 1:
            div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                          Inches(x1 - 0.005), Inches(MonthRow_y),
                                          Inches(0.010), Inches(MonthRow_h))
            solid_fill(div, RGBColor(
                min(theme["month_strip"][0] + 40, 255),
                min(theme["month_strip"][1] + 40, 255),
                min(theme["month_strip"][2] + 40, 255),
            ))

        add_textbox(slide, x0, MonthRow_y + 0.04, cell_w, MonthRow_h - 0.08,
                    m_label, font_size=9.5, bold=True,
                    color=theme["month_text"], align=PP_ALIGN.CENTER)

        vl = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                     Inches(x0), Inches(Grid_y0),
                                     Inches(0.009), Inches(Grid_y1 - Grid_y0))
        solid_fill(vl, theme["grid"])

    # Outer grid + label column borders
    gb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 Inches(Grid_x0), Inches(Grid_y0),
                                 Inches(Grid_w),  Inches(Grid_y1 - Grid_y0))
    gb.fill.background()
    set_shape_line(gb, theme["grid"], 0.75)

    lb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 Inches(M_left), Inches(Grid_y0),
                                 Inches(Left_w), Inches(Grid_y1 - Grid_y0))
    lb.fill.background()
    set_shape_line(lb, theme["grid"], 0.75)

    # ── Row geometry ──────────────────────────────────────────
    n_phase = sum(1 for r in rows if r.kind == "phase")
    n_task  = sum(1 for r in rows if r.kind in ("task", "milestone"))
    avail_h = Grid_y1 - Grid_y0
    base_content = n_phase * 0.33 + n_task * 0.265
    fill_scale   = min(avail_h / base_content, 1.45) if base_content else 1.0
    H_phase = 0.33  * fill_scale
    H_task  = 0.265 * fill_scale
    content_h = n_phase * H_phase + n_task * H_task
    if content_h > avail_h:
        s = avail_h / content_h
        H_phase *= s
        H_task  *= s

    row_y_mid: Dict[str, float] = {}

    y = Grid_y0
    for r in rows:
        h     = H_phase if r.kind == "phase" else H_task
        y_mid = y + h / 2

        if r.kind == "phase":
            band = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(M_left), Inches(y), Inches(Left_w), Inches(h))
            solid_fill(band, theme["phase_fill"])

            grid_band = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(Grid_x0), Inches(y), Inches(Grid_w), Inches(h))
            solid_fill(grid_band, theme["phase_fill"])

            add_textbox(slide, M_left + 0.07, y + 0.025, Left_w - 0.10, h - 0.05,
                        r.label, font_size=8.5, bold=True, color=theme["text"], vcenter=True)
        else:
            add_textbox(slide, M_left + 0.09, y + 0.015, Left_w - 0.16, h - 0.03,
                        r.label, font_size=7.5, color=theme["text"], vcenter=True)

            sep = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(M_left), Inches(y + h), Inches(Left_w + Grid_w), Inches(0.007))
            solid_fill(sep, theme["grid"])

            if r.row_id:
                row_y_mid[r.row_id] = y_mid
        y += h

    # ── Task bars ─────────────────────────────────────────────
    for r in rows:
        if r.kind != "task" or not r.start or not r.row_id:
            continue
        ym = row_y_mid.get(r.row_id)
        if ym is None:
            continue
        bh = H_task * 0.50
        x0 = x_of(r.start)
        w  = max(0.035, x_of(r.end + timedelta(days=1)) - x0)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x0), Inches(ym - bh / 2), Inches(w), Inches(bh))
        solid_fill(bar, theme["bar"])

    # ── Launch dashed line + label ────────────────────────────
    x_launch = x_of(milestone_dates["launch_to_market"])
    dashed_vline(slide, x_launch, Grid_y0, Grid_y1, theme["launch_line"])
    add_textbox(
        slide, x_launch + 0.04, Grid_y0 + 0.05, 1.40, 0.35,
        f"Launch to market\n{milestone_dates['launch_to_market'].strftime('%m/%d/%Y')}",
        font_size=7, bold=True, italic=True, color=theme["launch_line"],
    )

    # ── Milestone triangles ───────────────────────────────────
    for key in ["ioi_due", "loi_due", "close"]:
        if key not in milestone_dates:
            continue
        d   = milestone_dates[key]
        xm  = x_of(d)
        rid = milestone_row_map.get(key)
        ym  = row_y_mid.get(rid, Grid_y0 + 0.5) if rid else Grid_y0 + 0.5
        tw, th = 0.14, 0.12
        tri = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
            Inches(xm - tw / 2), Inches(ym - th - 0.03), Inches(tw), Inches(th))
        solid_fill(tri, theme["milestone"])
        add_textbox(
            slide, xm - 0.42, ym - 0.01, 0.88, 0.18,
            d.strftime("%m/%d/%Y"),
            font_size=7, bold=True,
            color=theme["muted_text"], align=PP_ALIGN.CENTER,
        )

    # ── Footer ────────────────────────────────────────────────
    script_dir    = os.path.dirname(os.path.abspath(__file__))
    logo_filename = theme["logo_file"]
    logo_path     = os.path.join(script_dir, logo_filename)
    if not os.path.exists(logo_path):
        alt = "logo_light.png" if logo_filename == "logo_dark.png" else "logo_dark.png"
        logo_path = os.path.join(script_dir, alt)

    logo_h        = 0.20
    footer_font   = 7.0
    footer_color  = theme["footer_text"]
    footer_text_x = M_left

    if os.path.exists(logo_path):
        try:
            from PIL import Image as PILImage
            with PILImage.open(logo_path) as img:
                iw, ih = img.size
            logo_w = logo_h * (iw / ih)
            slide.shapes.add_picture(
                logo_path,
                Inches(M_left),
                Inches(Footer_y + (Footer_h - logo_h) / 2),
                width=Inches(logo_w),
                height=Inches(logo_h),
            )
            footer_text_x = M_left + logo_w + 0.10
        except Exception:
            pass

    footer_note = "Note(s): Project plan is preliminary and subject to changes"
    sc_w = 1.60
    footer_center_w = SW - footer_text_x - sc_w - M_right - 0.05
    add_textbox(slide, footer_text_x, Footer_y + 0.04, footer_center_w, Footer_h - 0.04,
                footer_note, font_size=footer_font, color=footer_color)

    add_textbox(slide, SW - M_right - sc_w, Footer_y + 0.04, sc_w, Footer_h - 0.04,
                "Strictly Confidential",
                font_size=footer_font, color=footer_color, align=PP_ALIGN.RIGHT)

    prs.save(out_path)
    if isinstance(out_path, str):
        print(f"Wrote: {out_path}")


# ─────────────────────────────
# CLI
# ─────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--close",    required=True, help="Close date, e.g. 2026-11-07")
    ap.add_argument("--process",  required=True, choices=["standard", "accelerated"])
    ap.add_argument("--theme",    required=True, choices=["light", "dark"])
    ap.add_argument("--out",      default="timeline_slide.pptx")
    ap.add_argument("--subtitle", default="",
                    help="Bold subtitle line. Auto-generated from dates if omitted.")
    ap.add_argument("--top-label", default="[Insert section label here]",
                    dest="top_label",
                    help="Small label top-left, e.g. '3 | Process design and investor discussion'")
    args = ap.parse_args()

    render_timeline_slide(
        isoparse(args.close).date(),
        args.process,
        args.theme,
        args.out,
        subtitle=args.subtitle,
        top_label=args.top_label,
    )


if __name__ == "__main__":
    main()
